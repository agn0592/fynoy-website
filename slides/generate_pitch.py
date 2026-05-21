#!/usr/bin/env python3
import sys
import json
import os
import random
import datetime
import requests
from io import BytesIO

from dotenv import load_dotenv
from anthropic import Anthropic
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from google.auth.transport.requests import Request
from google.oauth2.credentials import Credentials
from google_auth_oauthlib.flow import InstalledAppFlow
from googleapiclient.discovery import build
from googleapiclient.http import MediaFileUpload

# ---------------------------------------------------------------------------
# Brand Color Palette (Fynoy Capital)
# ---------------------------------------------------------------------------
NAVY       = RGBColor(0x1E, 0x27, 0x61)   # #1E2761  title bars / dark slides
ICE_BLUE   = RGBColor(0xCA, 0xDC, 0xFC)   # #CADCFC  accent on DARK backgrounds only
ACCENT_BLUE = RGBColor(0xCA, 0xDC, 0xFC)  # #CADCFC
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK_NAVY  = RGBColor(0x0F, 0x16, 0x40)
RED_ACCENT = RGBColor(0xC0, 0x39, 0x2B)
LIGHT_GRAY = RGBColor(0xF4, 0xF6, 0xF9)

TITLE_FONT = "Georgia"
BODY_FONT  = "Calibri"
SLIDE_W    = Inches(13.333)
SLIDE_H    = Inches(7.5)

MARGIN      = Inches(0.5)
TITLE_BAR_H = Inches(1.0)
CONTENT_TOP = Inches(1.6)
BULLET_SPC  = Pt(10.8)


SYSTEM_PROMPT = """You are an elite quantitative analyst and presentation designer at Fynoy Capital, a top-tier hedge fund.
You receive raw JSON data from a form record containing assorted information about an investment opportunity.
Your task is to extract the most relevant, high-impact data and restructure it into English text for a 9-slide investment pitch presentation.

Use all provided input data as primary source. Search your knowledge for:
- The company's real Revenue Growth % and Operating Margin % for the Competitive Position table
- Real P/E ratios for the 3 competitors listed in input for the bar chart
- Real sector average EV/EBITDA for the Valuation slide
- Any recent material news about the company not in the input
Only include data you are confident is accurate. Do not fabricate numbers.

For the competitive position table, you MUST include the company's own data in the company row. Use real financial data: populate Revenue Growth and Operating Margin with actual figures. Do not leave the company row empty or with dashes.

IMPORTANT: The P/E, Forward P/E, and EV/EBITDA values in slide 6 metrics array are provided separately by the system from the original input. Do NOT override them. Only provide bar_chart and sector_avg_evebidta from your knowledge.

Output ONLY a valid JSON object matching the requested structure. Use short to medium professional bullet points (max 6 per slide).

Required 9 Slides (in order):
1. Cover (Company Name, Ticker, Sector, Date)
2. Company Overview (Business description, sector, country, business model, key revenue streams — 5 bullets)
3. Investment Thesis & Key Catalysts (Thesis statement 1-2 sentences, then exactly 4 catalyst descriptions)
4. Fundamental Analysis (Three categories: "Earnings Quality", "Business Model", "Governance" — each with 2-3 bullet points)
5. Competitive Position (3-4 market position bullets. Company + 3 competitors with Revenue Growth, Operating Margin, Market Position. Populate the company row with real data from your knowledge.)
6. Valuation (P/E, Forward P/E, EV/EBITDA as numbers. Bar chart data with company + 3 competitor P/E. Sector average EV/EBITDA number.)
7. Risks (Exactly 4 risks, each with a short title and a 1-sentence description)
8. Trade Setup (Direction, Entry Price, Stop Loss, Take Profit, Risk/Reward, Confidence Level as "Label: Value", plus holding_period)
9. Conclusion (Recommendation 1-2 sentences, score out of 10)

JSON Structure:
{
  "ticker": "AAPL",
  "company_name": "Apple Inc.",
  "sector": "Technology",
  "slides": [
    {"slide_index": 1, "title": "Cover", "content": []},
    {"slide_index": 2, "title": "Company Overview", "content": ["bullet 1", "bullet 2", "bullet 3", "bullet 4", "bullet 5"]},
    {"slide_index": 3, "title": "Investment Thesis & Key Catalysts", "thesis": "Thesis statement.", "catalysts": ["Catalyst 1", "Catalyst 2", "Catalyst 3", "Catalyst 4"]},
    {"slide_index": 4, "title": "Fundamental Analysis", "categories": [{"name": "Earnings Quality", "bullets": ["point 1", "point 2"]}, {"name": "Business Model", "bullets": ["point 1", "point 2"]}, {"name": "Governance", "bullets": ["point 1", "point 2"]}]},
    {"slide_index": 5, "title": "Competitive Position", "content": ["Market position bullet", "MOAT bullet", "Advantage bullet"], "company_metrics": {"revenue_growth": "8.5%", "operating_margin": "22.3%", "market_position": "Leader"}, "competitors": [{"name": "Comp 1", "revenue_growth": "6.2%", "operating_margin": "18.1%", "market_position": "Challenger"}, {"name": "Comp 2", "revenue_growth": "12.0%", "operating_margin": "15.5%", "market_position": "Niche"}, {"name": "Comp 3", "revenue_growth": "4.8%", "operating_margin": "20.1%", "market_position": "Follower"}]},
    {"slide_index": 6, "title": "Valuation", "metrics": [{"label": "P/E", "value": "18.5x"}, {"label": "Forward P/E", "value": "15.2x"}, {"label": "EV/EBITDA", "value": "12.3x"}], "bar_chart": [{"name": "Company", "pe": 18.5}, {"name": "Comp 1", "pe": 20.1}, {"name": "Comp 2", "pe": 15.8}, {"name": "Comp 3", "pe": 22.3}], "sector_avg_evebidta": "14.2x"},
    {"slide_index": 7, "title": "Risks", "risks": [{"title": "Risk Title", "description": "Short description."}, {"title": "Risk 2", "description": "Description."}, {"title": "Risk 3", "description": "Description."}, {"title": "Risk 4", "description": "Description."}]},
    {"slide_index": 8, "title": "Trade Setup", "content": ["Direction: Long", "Entry Price: $150", "Stop Loss: $135", "Take Profit: $185", "Risk/Reward: 2.3:1", "Confidence: High"], "holding_period": "6-12 months"},
    {"slide_index": 9, "title": "Conclusion", "content": ["Recommendation sentence.", "Score: 8/10"]}
  ]
}

IMPORTANT rules:
- Slide 4: use "categories" array with "name" and "bullets" for each of the 3 categories
- Slide 5: "company_metrics" object with revenue_growth, operating_margin, market_position. "competitors" array of 3.
- Slide 6: "metrics" array of 3 objects. "bar_chart" array of 4. "sector_avg_evebidta" string.
- Slide 7: "risks" array of exactly 4 objects with "title" and "description"
- Slide 9: last content item must start with "Score:"
- Ensure "slides" array has exactly 9 objects.
"""


def get_pexels_image(query: str, api_key: str, orientation: str = "landscape"):
    try:
        url = "https://api.pexels.com/v1/search"
        headers = {"Authorization": api_key}
        params = {"query": query, "per_page": 5, "orientation": orientation}
        resp = requests.get(url, headers=headers, params=params, timeout=10)
        if resp.status_code == 200:
            data = resp.json()
            if data.get("photos"):
                photo = random.choice(data["photos"])
                img_url = photo["src"]["large2x"]
                img_resp = requests.get(img_url, timeout=10)
                if img_resp.status_code == 200:
                    return BytesIO(img_resp.content)
    except Exception as e:
        print(f"  ⚠ Pexels fetch failed for '{query}': {e}")
    return None


def generate_content_with_claude(raw_json: dict) -> dict:
    api_key = os.environ.get("ANTHROPIC_API_KEY")
    if not api_key:
        raise ValueError("ANTHROPIC_API_KEY not found in environment variables.")
    client = Anthropic(api_key=api_key)
    response = client.messages.create(
        model="claude-sonnet-4-20250514",
        max_tokens=4500,
        system=SYSTEM_PROMPT,
        messages=[{"role": "user", "content": f"Generate the structured JSON for the 9-slide presentation based on this data:\n\n{json.dumps(raw_json)}"}],
        temperature=0.2,
    )
    text = response.content[0].text
    if "```json" in text:
        text = text.split("```json")[1].split("```")[0].strip()
    elif "```" in text:
        text = text.split("```")[1].strip()
    return json.loads(text)


def _add_bg(slide, color=LIGHT_GRAY):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def _add_rect(slide, left, top, width, height, fill=DARK_NAVY, border=ICE_BLUE, bw=Pt(1.5)):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border
        s.line.width = bw
    else:
        s.line.fill.background()
    return s


def _add_line(slide, left, top, width, height=Inches(0.025), color=ICE_BLUE):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def _run(p, text, font=BODY_FONT, size=Pt(13), color=DARK_NAVY,
         bold=False, italic=False, align=None):
    p.text = ""
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = size
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    if align is not None:
        p.alignment = align


def _title_bar(slide, text):
    _add_rect(slide, Inches(0), Inches(0), SLIDE_W, TITLE_BAR_H, fill=NAVY, border=None)
    tb = slide.shapes.add_textbox(MARGIN, Inches(0.27), Inches(12), Inches(0.45))
    _run(tb.text_frame.paragraphs[0], text,
         font=TITLE_FONT, size=Pt(36), color=WHITE, bold=True)


def _bullets(slide, items, left, top, width, height, size=Pt(13),
             text_color=DARK_NAVY, char="•", char_color=ACCENT_BLUE, max_n=6):
    items = [str(i) for i in (items or []) if i][:max_n]
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        rb = p.add_run()
        rb.text = f"{char}  "
        rb.font.name = BODY_FONT
        rb.font.size = size
        rb.font.color.rgb = char_color
        rb.font.bold = True
        rt = p.add_run()
        rt.text = item
        rt.font.name = BODY_FONT
        rt.font.size = size
        rt.font.color.rgb = text_color
        p.space_before = Pt(4)
        p.space_after = Pt(16)
    return tb


def _watermark(slide, cover=False, dark_slide=False, shift_down=False):
    y = Inches(7.5 - 0.15 - 0.25)
    if cover:
        y = y - Inches(0.3)
    if shift_down:
        y = y + Inches(0.12)
    x = Inches(13.333 - 0.2 - 4.3)
    tb = slide.shapes.add_textbox(x, y, Inches(4.3), Inches(0.25))
    wm_color = ICE_BLUE if dark_slide else DARK_NAVY
    _run(tb.text_frame.paragraphs[0], "© Fynoy Capital — Confidential",
         font=BODY_FONT, size=Pt(9), color=wm_color, align=PP_ALIGN.RIGHT)


def _safe_img(stream, slide, left, top, width=None, height=None):
    if stream is None:
        return None
    try:
        stream.seek(0)
        kw = {}
        if width:
            kw["width"] = width
        if height:
            kw["height"] = height
        return slide.shapes.add_picture(stream, left, top, **kw)
    except Exception as e:
        print(f"  ⚠ Image insert failed: {e}")
        return None


def _slide_cover(prs, data, content, images):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    img = images.get("cover")
    if img:
        _safe_img(img, slide, Inches(0), Inches(0), width=SLIDE_W, height=SLIDE_H)
    _add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=NAVY, border=None)
    _add_line(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.05))
    name = content.get("company_name", "Company Name")
    tb = slide.shapes.add_textbox(Inches(1), Inches(2.0), Inches(11.333), Inches(1.2))
    _run(tb.text_frame.paragraphs[0], name,
         font=TITLE_FONT, size=Pt(52), color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    _add_line(slide, Inches(4.5), Inches(3.4), Inches(4.333), Inches(0.04))
    ticker = content.get("ticker", "")
    sector = content.get("sector", "")
    now = datetime.datetime.now()
    sub = "  ·  ".join(p for p in [ticker, sector, now.strftime("%B %d, %Y")] if p)
    stb = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(11.333), Inches(0.5))
    _run(stb.text_frame.paragraphs[0], sub, size=Pt(16), color=ICE_BLUE, align=PP_ALIGN.CENTER)
    valid = now + datetime.timedelta(days=30)
    vtb = slide.shapes.add_textbox(Inches(1), Inches(4.4), Inches(11.333), Inches(0.4))
    _run(vtb.text_frame.paragraphs[0], f"Valid until: {valid.strftime('%B %d, %Y')}",
         size=Pt(12), color=ICE_BLUE, align=PP_ALIGN.CENTER)
    _add_line(slide, Inches(0), Inches(7.15), SLIDE_W, Inches(0.05))
    _watermark(slide, cover=True, dark_slide=True)


def _slide_overview(prs, data, content, images):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, LIGHT_GRAY)
    _title_bar(slide, data.get("title", "Company Overview"))
    _bullets(slide, data.get("content", []), MARGIN, CONTENT_TOP, Inches(6.83), Inches(5.5))
    img_l = Inches(7.33)
    img_w = Inches(5.5)
    img_h = SLIDE_H - TITLE_BAR_H
    img = images.get("overview")
    if img:
        _safe_img(img, slide, img_l, TITLE_BAR_H, width=img_w, height=img_h)
    else:
        _add_rect(slide, img_l, TITLE_BAR_H, img_w, img_h, fill=WHITE, border=ICE_BLUE, bw=Pt(2))
        ntb = slide.shapes.add_textbox(img_l + Inches(0.2), TITLE_BAR_H + Inches(2.0), Inches(5.1), Inches(1.0))
        _run(ntb.text_frame.paragraphs[0], content.get("company_name", ""), font=TITLE_FONT, size=Pt(28), color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    _watermark(slide, dark_slide=False)


def _slide_thesis(prs, data, content, images):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, LIGHT_GRAY)
    _title_bar(slide, "Investment Thesis")
    thesis = data.get("thesis", "")
    catalysts = data.get("catalysts", data.get("content", []))[:4]
    ttb = slide.shapes.add_textbox(MARGIN, CONTENT_TOP, Inches(12.3), Inches(0.9))
    ttb.text_frame.word_wrap = True
    _run(ttb.text_frame.paragraphs[0], str(thesis), size=Pt(15), color=DARK_NAVY, italic=True)
    _add_line(slide, MARGIN, Inches(2.3), Inches(12.3), Inches(0.025))
    sh = slide.shapes.add_textbox(MARGIN, Inches(2.5), Inches(5), Inches(0.4))
    _run(sh.text_frame.paragraphs[0], "Key Catalysts", size=Pt(16), color=ACCENT_BLUE, bold=True)
    bw = Inches(5.9)
    bh = Inches(1.55)
    xs = [MARGIN, Inches(6.7)]
    ys = [Inches(3.1), Inches(4.85)]
    for i, cat in enumerate(catalysts):
        x = xs[i % 2]
        y = ys[i // 2]
        _add_rect(slide, x, y, bw, bh, fill=NAVY, border=None)
        ntb = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.2), Inches(0.5), Inches(0.4))
        _run(ntb.text_frame.paragraphs[0], str(i + 1), size=Pt(22), color=ICE_BLUE, bold=True)
        ctb = slide.shapes.add_textbox(x + Inches(0.75), y + Inches(0.2), bw - Inches(1.0), bh - Inches(0.4))
        ctb.text_frame.word_wrap = True
        _run(ctb.text_frame.paragraphs[0], f"→ {cat}", size=Pt(13), color=WHITE)
    _watermark(slide, dark_slide=False)


def _slide_fundamental(prs, data, content, images):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, LIGHT_GRAY)
    _title_bar(slide, data.get("title", "Fundamental Analysis"))
    categories = data.get("categories", [])
    if not categories:
        for b in data.get("content", [])[:3]:
            parts = str(b).split(":", 1)
            if len(parts) == 2:
                categories.append({"name": parts[0].strip(), "bullets": [parts[1].strip()]})
            else:
                categories.append({"name": "Analysis", "bullets": [str(b)]})
    headers = ["Earnings Quality", "Business Model", "Governance"]
    while len(categories) < 3:
        categories.append({"name": headers[len(categories)], "bullets": []})
    cw = Inches(3.9)
    ch = Inches(5.6)
    gap = Inches(0.35)
    for i, cat in enumerate(categories[:3]):
        x = MARGIN + (cw + gap) * i
        y = CONTENT_TOP
        _add_rect(slide, x, y, cw, ch, fill=WHITE, border=ACCENT_BLUE, bw=Pt(1.5))
        _add_rect(slide, x, y, cw, Inches(0.5), fill=ACCENT_BLUE, border=None)
        htb = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.05), cw - Inches(0.2), Inches(0.4))
        _run(htb.text_frame.paragraphs[0], cat.get("name", headers[i]), size=Pt(14), color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        _bullets(slide, cat.get("bullets", []), x + Inches(0.15), y + Inches(0.7), cw - Inches(0.3), ch - Inches(0.9), size=Pt(13), max_n=4)
    _watermark(slide, dark_slide=False, shift_down=True)


def _slide_competitive(prs, data, content, images):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, LIGHT_GRAY)
    _title_bar(slide, data.get("title", "Competitive Position"))
    bull = data.get("content", [])
    competitors = data.get("competitors", [])
    company_name = content.get("company_name", "Company")
    cm = data.get("company_metrics", {})
    _bullets(slide, bull[:4], MARGIN, CONTENT_TOP, Inches(6.5), Inches(3.2), char="—", char_color=ACCENT_BLUE, size=Pt(13))
    tl = Inches(7.3)
    col_w = [Inches(1.55), Inches(1.3), Inches(1.3), Inches(1.35)]
    total_w = sum(c for c in col_w)
    rh = Inches(0.55)
    tt = CONTENT_TOP
    cols = ["", "Revenue Growth", "Operating Margin", "Market Position"]
    _add_rect(slide, tl, tt, total_w, rh, fill=NAVY, border=None)
    cx = tl
    for ci, h in enumerate(cols):
        w = col_w[ci]
        htb = slide.shapes.add_textbox(cx + Inches(0.04), tt + Inches(0.06), w - Inches(0.08), rh - Inches(0.12))
        _run(htb.text_frame.paragraphs[0], h, size=Pt(10), color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        cx += w
    rows = [{"name": company_name, "rg": cm.get("revenue_growth", "—"), "om": cm.get("operating_margin", "—"), "mp": cm.get("market_position", "—")}]
    for c in competitors[:3]:
        rows.append({"name": c.get("name", "—"), "rg": c.get("revenue_growth", "—"), "om": c.get("operating_margin", "—"), "mp": c.get("market_position", "—")})
    while len(rows) < 4:
        rows.append({"name": "—", "rg": "—", "om": "—", "mp": "—"})
    for ri, rd in enumerate(rows):
        ry = tt + rh * (ri + 1)
        bg = ICE_BLUE if ri == 0 else (WHITE if ri % 2 == 1 else LIGHT_GRAY)
        tc = NAVY if ri == 0 else DARK_NAVY
        _add_rect(slide, tl, ry, total_w, rh, fill=bg, border=ICE_BLUE, bw=Pt(0.5))
        vals = [rd["name"], rd["rg"], rd["om"], rd["mp"]]
        cx = tl
        for ci, v in enumerate(vals):
            w = col_w[ci]
            ctb = slide.shapes.add_textbox(cx + Inches(0.04), ry + Inches(0.06), w - Inches(0.08), rh - Inches(0.12))
            _run(ctb.text_frame.paragraphs[0], str(v), size=Pt(10), color=tc, bold=(ci == 0 or ri == 0), align=PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT)
            cx += w
    _watermark(slide, dark_slide=False)


def _slide_valuation(prs, data, content, images):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, LIGHT_GRAY)
    _title_bar(slide, data.get("title", "Valuation"))
    protected = content.get("_protected_valuation", {})
    metrics = [
        {"label": "P/E",        "value": protected.get("pe", data.get("metrics", [{}])[0].get("value", "N/A") if data.get("metrics") else "N/A")},
        {"label": "Forward P/E", "value": protected.get("fwd_pe", data.get("metrics", [{},{}])[1].get("value", "N/A") if len(data.get("metrics",[])) > 1 else "N/A")},
        {"label": "EV/EBITDA",  "value": protected.get("ev_ebitda", data.get("metrics", [{},{},{}])[2].get("value", "N/A") if len(data.get("metrics",[])) > 2 else "N/A")},
    ]
    bar_chart = data.get("bar_chart", [])
    sector_avg = data.get("sector_avg_evebidta", "")
    bw = Inches(3.8)
    bh = Inches(1.6)
    gap = Inches(0.35)
    for i, m in enumerate(metrics[:3]):
        x = MARGIN + (bw + gap) * i
        y = CONTENT_TOP
        _add_rect(slide, x, y, bw, bh, fill=WHITE, border=NAVY, bw=Pt(2))
        ltb = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.1), bw - Inches(0.3), Inches(0.3))
        _run(ltb.text_frame.paragraphs[0], m["label"], size=Pt(12), color=ACCENT_BLUE, align=PP_ALIGN.CENTER)
        vtb = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.4), bw - Inches(0.3), Inches(1.0))
        _run(vtb.text_frame.paragraphs[0], m["value"], size=Pt(44), color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    ct = Inches(3.6)
    sh = slide.shapes.add_textbox(MARGIN, ct, Inches(6), Inches(0.35))
    _run(sh.text_frame.paragraphs[0], "P/E Ratio Comparison", size=Pt(16), color=NAVY, bold=True)
    if bar_chart:
        max_pe = max((float(b.get("pe", 0)) for b in bar_chart), default=1) or 1
        bar_h = Inches(0.45)
        label_w = Inches(2.5)
        max_bw = Inches(8.0)
        by = ct + Inches(0.5)
        for i, bar in enumerate(bar_chart[:4]):
            name = bar.get("name", "—")
            pe = float(bar.get("pe", 0))
            y = by + (bar_h + Inches(0.2)) * i
            is_co = (i == 0)
            lt = slide.shapes.add_textbox(MARGIN, y, label_w, bar_h)
            _run(lt.text_frame.paragraphs[0], name, size=Pt(11), color=NAVY, bold=is_co, align=PP_ALIGN.RIGHT)
            w = max_bw * (pe / max_pe) if max_pe > 0 else Inches(0.5)
            bc = NAVY if is_co else ICE_BLUE
            _add_rect(slide, MARGIN + label_w + Inches(0.15), y + Inches(0.04), w, bar_h - Inches(0.08), fill=bc, border=None)
            vt = slide.shapes.add_textbox(MARGIN + label_w + Inches(0.25) + w, y + Inches(0.04), Inches(0.8), bar_h - Inches(0.08))
            _run(vt.text_frame.paragraphs[0], f"{pe:.1f}x", size=Pt(11), color=NAVY, bold=True)
    if sector_avg:
        sat = slide.shapes.add_textbox(MARGIN, Inches(6.5), Inches(12), Inches(0.35))
        _run(sat.text_frame.paragraphs[0], f"◆  Sector Average EV/EBITDA: {sector_avg}", size=Pt(13), color=ACCENT_BLUE, bold=True)
    _watermark(slide, dark_slide=False)


def _slide_risks(prs, data, content, images):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, LIGHT_GRAY)
    _title_bar(slide, data.get("title", "Risks"))
    risks = data.get("risks", [])[:4]
    while len(risks) < 4:
        risks.append({"title": "—", "description": ""})
    cw = Inches(5.6)
    ch = Inches(2.2)
    gap_x = Inches(0.5)
    gap_y = Inches(0.4)
    total_w = cw * 2 + gap_x
    start_x = (SLIDE_W - total_w) / 2
    start_y = Inches(1.8)
    xs = [start_x, start_x + cw + gap_x]
    ys = [start_y, start_y + ch + gap_y]
    for i, risk in enumerate(risks):
        col = i % 2
        row = i // 2
        x = xs[col]
        y = ys[row]
        _add_rect(slide, x, y, cw, ch, fill=WHITE, border=ICE_BLUE, bw=Pt(1))
        _add_rect(slide, x, y, Inches(0.06), ch, fill=RED_ACCENT, border=None)
        ttb = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.15), cw - Inches(0.35), Inches(0.4))
        p = ttb.text_frame.paragraphs[0]
        rw = p.add_run()
        rw.text = "⚠  "
        rw.font.name = BODY_FONT
        rw.font.size = Pt(13)
        rw.font.color.rgb = RED_ACCENT
        rw.font.bold = True
        rt = p.add_run()
        rt.text = risk.get("title", "—")
        rt.font.name = BODY_FONT
        rt.font.size = Pt(13)
        rt.font.color.rgb = DARK_NAVY
        rt.font.bold = True
        desc = risk.get("description", "")
        if desc:
            dtb = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.65), cw - Inches(0.35), Inches(1.3))
            dtb.text_frame.word_wrap = True
            _run(dtb.text_frame.paragraphs[0], desc, size=Pt(12), color=DARK_NAVY)
    _watermark(slide, dark_slide=False)


def _slide_trade(prs, data, content, images):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, LIGHT_GRAY)
    _title_bar(slide, data.get("title", "Trade Setup"))
    bullets = data.get("content", [])
    holding = data.get("holding_period", "")
    pairs = []
    for b in bullets:
        s = str(b)
        if ":" in s:
            parts = s.split(":", 1)
            pairs.append((parts[0].strip(), parts[1].strip()))
        else:
            pairs.append(("Info", s))
    symbols = {"Direction": "↑", "Entry Price": "$", "Stop Loss": "$", "Take Profit": "$", "Risk/Reward": "R/R", "Confidence": "✓"}
    bw = Inches(3.8)
    bh = Inches(2.0)
    gx = Inches(0.45)
    gy = Inches(0.35)
    y1 = Inches(1.8)
    y2 = y1 + bh + gy
    total_w = bw * 3 + gx * 2
    start_x = (SLIDE_W - total_w) / 2
    for i, (label, value) in enumerate(pairs[:6]):
        col = i % 3
        row = i // 3
        x = start_x + (bw + gx) * col
        y = y1 if row == 0 else y2
        sym = symbols.get(label, "●")
        _add_rect(slide, x, y, bw, bh, fill=WHITE, border=NAVY, bw=Pt(2))
        ltb = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.15), bw - Inches(0.2), Inches(0.3))
        p = ltb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        rs = p.add_run()
        rs.text = f"{sym}  "
        rs.font.name = BODY_FONT
        rs.font.size = Pt(11)
        rs.font.color.rgb = ACCENT_BLUE
        rs.font.bold = True
        rl = p.add_run()
        rl.text = label
        rl.font.name = BODY_FONT
        rl.font.size = Pt(11)
        rl.font.color.rgb = ACCENT_BLUE
        rl.font.bold = True
        vtb = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.55), bw - Inches(0.2), Inches(1.2))
        _run(vtb.text_frame.paragraphs[0], value, size=Pt(28), color=DARK_NAVY, bold=True, align=PP_ALIGN.CENTER)
    if holding:
        ht = slide.shapes.add_textbox(MARGIN, Inches(5.8), Inches(12.3), Inches(0.4))
        _run(ht.text_frame.paragraphs[0], f"Expected Holding Period: {holding}", size=Pt(13), color=DARK_NAVY, align=PP_ALIGN.CENTER)
    _watermark(slide, dark_slide=False)


def _slide_conclusion(prs, data, content, images):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, NAVY)
    bullets = data.get("content", [])
    ttb = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(11.333), Inches(0.8))
    _run(ttb.text_frame.paragraphs[0], data.get("title", "Conclusion"),
         font=TITLE_FONT, size=Pt(36), color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    _add_line(slide, Inches(4.5), Inches(1.7), Inches(4.333), Inches(0.03))
    rec = bullets[0] if bullets else "Final Recommendation"
    rtb = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.333), Inches(1.4))
    rtb.text_frame.word_wrap = True
    _run(rtb.text_frame.paragraphs[0], str(rec), font=TITLE_FONT, size=Pt(22), color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    score_num = ""
    for b in bullets:
        s = str(b)
        if s.lower().startswith("score"):
            parts = s.split(":", 1)
            raw = parts[1].strip() if len(parts) == 2 else s
            raw = raw.replace("/10", "").strip()
            score_num = raw
            break
    if not score_num:
        score_num = "—"
    stb = slide.shapes.add_textbox(Inches(3.5), Inches(3.8), Inches(6.333), Inches(1.4))
    p = stb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    rs = p.add_run()
    rs.text = f"{score_num}/10"
    rs.font.name = BODY_FONT
    rs.font.size = Pt(64)
    rs.font.color.rgb = ICE_BLUE
    rs.font.bold = True
    ltb = slide.shapes.add_textbox(Inches(3.5), Inches(5.2), Inches(6.333), Inches(0.4))
    _run(ltb.text_frame.paragraphs[0], "Investment Score", size=Pt(14), color=ICE_BLUE, align=PP_ALIGN.CENTER)
    tag = slide.shapes.add_textbox(Inches(3.5), Inches(6.2), Inches(6.333), Inches(0.35))
    _run(tag.text_frame.paragraphs[0], "Prepared by Fynoy Capital", size=Pt(11), color=ICE_BLUE, align=PP_ALIGN.CENTER)
    _add_line(slide, Inches(0), Inches(7.15), SLIDE_W, Inches(0.05))
    _watermark(slide, dark_slide=True, shift_down=True)


SLIDE_BUILDERS = [
    _slide_cover,
    _slide_overview,
    _slide_thesis,
    _slide_fundamental,
    _slide_competitive,
    _slide_valuation,
    _slide_risks,
    _slide_trade,
    _slide_conclusion,
]


def build_presentation(content: dict, images: dict) -> str:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slides_data = content.get("slides", [])
    for idx, builder in enumerate(SLIDE_BUILDERS):
        d = slides_data[idx] if idx < len(slides_data) else {"title": "", "content": []}
        builder(prs, d, content, images)
    ts = datetime.datetime.now().strftime("%Y-%m-%d")
    ticker = content.get("ticker", "UNKNOWN")
    fn = f"{ticker}_{ts}.pptx"
    prs.save(fn)
    return fn


def get_or_create_folder(service, folder_name, parent_id=None):
    query = f"name='{folder_name}' and mimeType='application/vnd.google-apps.folder' and trashed=false"
    if parent_id:
        query += f" and '{parent_id}' in parents"
    results = service.files().list(q=query, spaces='drive', fields='files(id, name)').execute()
    items = results.get('files', [])
    if not items:
        meta = {'name': folder_name, 'mimeType': 'application/vnd.google-apps.folder'}
        if parent_id:
            meta['parents'] = [parent_id]
        folder = service.files().create(body=meta, fields='id').execute()
        return folder.get('id')
    return items[0].get('id')


def upload_to_drive(filename: str):
    creds_path = os.environ.get("GOOGLE_DRIVE_CREDENTIALS_PATH", "credentials.json")
    token_path = 'token.json'
    scopes = ['https://www.googleapis.com/auth/drive.file']
    creds = None
    if os.path.exists(token_path):
        creds = Credentials.from_authorized_user_file(token_path, scopes)
    if not creds or not creds.valid:
        if creds and creds.expired and creds.refresh_token:
            try:
                creds.refresh(Request())
            except Exception as e:
                print(f"Token refresh failed: {e}. Will re-authenticate.")
                creds = None
        if not creds:
            if not os.path.exists(creds_path):
                print(f"Warning: Google Drive OAuth credentials not found at '{creds_path}'. Skipping upload.")
                return
            flow = InstalledAppFlow.from_client_secrets_file(creds_path, scopes)
            creds = flow.run_local_server(port=0)
        with open(token_path, 'w') as token:
            token.write(creds.to_json())
    try:
        service = build('drive', 'v3', credentials=creds)
        q = "name='Fynoy Capital' and mimeType='application/vnd.google-apps.folder' and trashed=false"
        results = service.files().list(q=q, spaces='drive', fields='files(id, name)').execute()
        items = results.get('files', [])
        if items:
            fid = items[0].get('id')
            pid = get_or_create_folder(service, 'Proposals', fid)
        else:
            fid = get_or_create_folder(service, 'Fynoy Capital')
            pid = get_or_create_folder(service, 'Proposals', fid)
        meta = {'name': filename, 'parents': [pid]}
        media = MediaFileUpload(filename, mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation')
        uf = service.files().create(body=meta, media_body=media, fields='id').execute()
        print(f"File uploaded. ID: {uf.get('id')}")
    except Exception as e:
        print(f"Error during Google Drive upload: {e}")


def main():
    load_dotenv()
    if len(sys.argv) > 1:
        if os.path.exists(sys.argv[1]):
            with open(sys.argv[1], 'r', encoding='utf-8') as f:
                raw_data_str = f.read()
        else:
            raw_data_str = sys.argv[1]
    else:
        raw_data_str = sys.stdin.read()
    try:
        raw_json = json.loads(raw_data_str)
    except json.JSONDecodeError:
        print("Error: Input is not valid JSON.", file=sys.stderr)
        sys.exit(1)
    if not raw_json:
        print("Error: Input JSON is empty.", file=sys.stderr)
        sys.exit(1)
    sc = generate_content_with_claude(raw_json)
    company = sc.get("company_name", "")
    sector = sc.get("sector", "")
    pkey = os.environ.get("PEXELS_API_KEY")
    protected_val = {}
    raw_fields = raw_json.get("fields", raw_json)
    for key, canon in [("P/E", "pe"), ("PE", "pe"), ("pe", "pe"), ("p_e", "pe"),
                       ("Forward P/E", "fwd_pe"), ("forward_pe", "fwd_pe"), ("fwd_pe", "fwd_pe"),
                       ("EV/EBITDA", "ev_ebitda"), ("ev_ebitda", "ev_ebitda")]:
        v = raw_fields.get(key)
        if v and canon not in protected_val:
            protected_val[canon] = str(v)
    sc["_protected_valuation"] = protected_val
    imgs = {}
    imgs["cover"] = get_pexels_image(f"{company} {sector} corporate header", pkey, orientation="landscape") if pkey else None
    imgs["overview"] = get_pexels_image(f"{company} {sector} corporate business operations", pkey, orientation="portrait") if pkey else None
    fn = build_presentation(sc, imgs)
    print(f"Saved: {fn}")
    upload_to_drive(fn)


if __name__ == "__main__":
    main()
